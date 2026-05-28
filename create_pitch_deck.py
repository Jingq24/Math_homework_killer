# -*- coding: utf-8 -*-
"""Generate fundraising pitch deck for 高数作业杀手 (Math Homework Killer)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)         # Deep navy
BG_CARD = RGBColor(0x25, 0x25, 0x3A)          # Card background
ACCENT = RGBColor(0x00, 0xD2, 0xFF)           # Cyan accent
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)          # Purple accent
GREEN = RGBColor(0x10, 0xB9, 0x81)            # Success green
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)           # Warning orange
RED = RGBColor(0xEF, 0x44, 0x44)              # Red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
DARK_TEXT = RGBColor(0xE2, 0xE8, 0xF0)

# ── Helper functions ───────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """Add a text box with single style."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, line_spacing=1.5, font_name="Microsoft YaHei"):
    """Add textbox with multiple styled lines. Each line is (text, bold, size_override, color_override)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, size, clr = line_data, False, font_size, color
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            size = line_data[2] if len(line_data) > 2 else font_size
            clr = line_data[3] if len(line_data) > 3 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return tf

def add_card(slide, left, top, width, height, color=BG_CARD):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_metric_card(slide, left, top, number, label, num_color=ACCENT):
    """Add a KPI metric card."""
    add_card(slide, left, top, 2.8, 1.5)
    add_textbox(slide, left + 0.2, top + 0.15, 2.4, 0.7, number,
                font_size=36, color=num_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.2, top + 0.9, 2.4, 0.4, label,
                font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

def new_slide():
    """Create a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    return slide

def add_section_title(slide, title, subtitle=None):
    """Add consistent section header."""
    add_accent_line(slide, 0.8, 0.7, 0.6)
    add_textbox(slide, 0.8, 0.85, 11, 0.7, title, font_size=32, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 1.45, 11, 0.5, subtitle, font_size=14, color=GRAY)

def add_page_number(slide, num):
    """Add page number."""
    add_textbox(slide, 12.3, 7.05, 0.8, 0.35, str(num),
                font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ────────────────────────────────────────────────────────────────
slide = new_slide()
# Large accent bar
add_accent_line(slide, 2.5, 2.6, 8.3, ACCENT)
add_textbox(slide, 2.5, 1.5, 8.3, 1.0, "高数作业杀手",
            font_size=56, bold=True)
add_textbox(slide, 2.5, 2.85, 8.3, 0.6, "AI-Powered Math Homework Assistant · WeChat Mini-Program",
            font_size=18, color=ACCENT)
add_textbox(slide, 2.5, 3.7, 8.3, 1.2,
            "拍照上传 → AI 秒解 → 手写体输出\n让高等数学作业不再头疼",
            font_size=20, color=GRAY)
add_textbox(slide, 2.5, 5.5, 8.3, 0.5, "融资路演 · 2026",
            font_size=14, color=LIGHT_GRAY)
add_page_number(slide, 1)

# ────────────────────────────────────────────────────────────────
# SLIDE 2: The Problem
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "市场痛点", "THE PROBLEM")

problems = [
    ("📐", "高数作业耗时长", "一道证明题平均耗时30-60分钟，\n理工科学生每周作业量超10题"),
    ("😰", "解题资源匮乏", "教材答案只有最终结果无过程，\n网课答疑排队数小时，效率极低"),
    ("📝", "手写体抄写麻烦", "找到答案后仍需手工抄写，\nAI对话输出的LaTeX格式无法直接使用"),
    ("💰", "现有工具不匹配", "Photomath不支持中文题目，\n国内暂无高数解题+手写体一站式工具"),
]

for i, (icon, title, desc) in enumerate(problems):
    x = 0.8 + i * 3.1
    add_card(slide, x, 2.3, 2.8, 3.8)
    add_textbox(slide, x + 0.3, 2.5, 2.2, 0.6, f"{icon}  {title}",
                font_size=18, bold=True)
    add_textbox(slide, x + 0.3, 3.3, 2.2, 2.5, desc,
                font_size=13, color=GRAY)

# Stats
add_metric_card(slide, 0.8, 6.4, "1200万+", "中国理工科大学生")
add_metric_card(slide, 3.9, 6.4, "3~5h/周", "平均高数作业耗时")
add_metric_card(slide, 7.0, 6.4, "¥50亿+", "大学生教育辅助市场")
add_metric_card(slide, 10.1, 6.4, "0个", "国内同类竞品")
add_page_number(slide, 2)

# ────────────────────────────────────────────────────────────────
# SLIDE 3: Our Solution
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "产品方案", "OUR SOLUTION")

# Flow diagram
steps = [
    ("📷", "拍照上传", "手机拍照\n或相册选取"),
    ("🧠", "AI 解题", "视觉大模型\n识别+推理"),
    ("✍️", "手写体渲染", "11款中文手写字体\n真实笔迹效果"),
    ("📄", "一键导出", "PDF下载/打印\n多题合并"),
]
for i, (icon, title, desc) in enumerate(steps):
    x = 0.8 + i * 3.1
    add_card(slide, x, 2.3, 2.8, 3.0)
    add_textbox(slide, x + 0.3, 2.5, 2.2, 0.5, icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.1, 2.2, 0.5, title, font_size=20, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.7, 2.2, 1.0, desc, font_size=14,
                color=GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow between cards (except last)
    if i < 3:
        add_textbox(slide, x + 2.85, 3.3, 0.3, 0.5, "→", font_size=28,
                    color=ACCENT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 5.8, 11.7, 1.2,
            "核心体验：微信小程序搜索「高数作业杀手」→ 拍题 → 等待10-20秒 → 获得手写体解答 PDF → 打印或保存\n"
            "支持高等数学、线性代数、概率统计等全部大学数学课程，覆盖中英文混合题目",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 3)

# ────────────────────────────────────────────────────────────────
# SLIDE 4: Tech Architecture
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "技术架构", "TECHNICAL ARCHITECTURE")

# Two columns
add_card(slide, 0.8, 2.3, 5.8, 4.5)
add_textbox(slide, 1.1, 2.5, 5.3, 0.5, "后端架构 (FastAPI + Celery)", font_size=18, bold=True, color=ACCENT)
backend_lines = [
    ("AI 引擎层", True, 14, ACCENT2),
    ("  · 兼容 OpenAI API 的视觉大模型调用", False, 12, GRAY),
    ("  · 支持 DashScope / SiliconFlow / Zhipu / Ollama", False, 12, GRAY),
    ("  · 智能后处理：去重、排版规范化", False, 12, GRAY),
    ("", False, 8, WHITE),
    ("渲染引擎层 (替代 Word COM)", True, 14, ACCENT2),
    ("  · LaTeX → HTML 语义转换", False, 12, GRAY),
    ("  · 字符级手写体抖动 (字号/位置/间距/字体切换)", False, 12, GRAY),
    ("  · 11款中文手写字体 + 3款数学字体", False, 12, GRAY),
    ("  · WeasyPrint / Playwright → PDF 输出", False, 12, GRAY),
    ("", False, 8, WHITE),
    ("基础设施", True, 14, ACCENT2),
    ("  · PostgreSQL + Redis + OSS 对象存储", False, 12, GRAY),
    ("  · Celery 异步任务队列 + Docker 容器化部署", False, 12, GRAY),
]
add_multiline(slide, 1.1, 3.0, 5.3, 3.5, backend_lines)

add_card(slide, 7.1, 2.3, 5.4, 4.5)
add_textbox(slide, 7.4, 2.5, 4.8, 0.5, "前端架构 (uni-app + Vue 3)", font_size=18, bold=True, color=ACCENT)
frontend_lines = [
    ("微信小程序端", True, 14, ACCENT2),
    ("  · uni-app 跨端框架，代码复用率高", False, 12, GRAY),
    ("  · 4个主Tab：首页/解题/记录/我的", False, 12, GRAY),
    ("  · 摄像头拍照 + 相册选图 + 图片压缩", False, 12, GRAY),
    ("", False, 8, WHITE),
    ("核心交互流程", True, 14, ACCENT2),
    ("  · 上传 → 轮询状态(3s) → 结果预览", False, 12, GRAY),
    ("  · Canvas 内嵌预览 + PDF WebView", False, 12, GRAY),
    ("  · 批量模式：多题一键求解+合并", False, 12, GRAY),
    ("  · 微信支付 V3 + 积分系统", False, 12, GRAY),
    ("", False, 8, WHITE),
    ("技术亮点", True, 14, ACCENT2),
    ("  · 分包加载，主包 < 2MB", False, 12, GRAY),
    ("  · 微信订阅消息通知", False, 12, GRAY),
    ("  · 激励广告组件预留", False, 12, GRAY),
]
add_multiline(slide, 7.4, 3.0, 4.8, 3.5, frontend_lines)

add_page_number(slide, 4)

# ────────────────────────────────────────────────────────────────
# SLIDE 5: Handwriting Engine (Key Differentiator)
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "核心壁垒：手写体渲染引擎", "CORE TECHNOLOGY MOAT")

add_card(slide, 0.8, 2.3, 5.8, 4.5)
add_textbox(slide, 1.1, 2.5, 5.3, 0.5, "11款中文手写字体库", font_size=18, bold=True, color=ACCENT)
fonts = [
    "美玉体 (MEIYUJW) — 默认字体",
    "世界那么大 (SJNMDJW)",
    "方正静蕾简体 (FZJingLeiS-R-GB)",
    "方正行楷 (AR GbxingkaiGB ExtraBold)",
    "汉仪晶润简 (HY JingBRJ)",
    "汉仪品体简 (HYPPTiJ)",
    "伯乐雅韵体 (BoLeYaYati)",
    "伯乐竹笋体 (BoLeZhuSunti)",
    "华康翩翩体 (DFPHanziPenW3-GB)",
    "伯乐俏皮体",
    "伯乐童年体",
]
for i, f in enumerate(fonts):
    add_textbox(slide, 1.3, 3.15 + i * 0.32, 5.0, 0.3, f"· {f}",
                font_size=12, color=GRAY)

add_card(slide, 7.1, 2.3, 5.4, 4.5)
add_textbox(slide, 7.4, 2.5, 4.8, 0.5, "随机扰动算法", font_size=18, bold=True, color=ACCENT)
jitter_lines = [
    ("字符级抖动", True, 14, ACCENT2),
    ("  · 字号随机抖动 ±12%", False, 12, GRAY),
    ("  · 基线位置偏移 ±5pt × 0.35", False, 12, GRAY),
    ("  · 字符间距波动 ±5pt", False, 12, GRAY),
    ("  · 横向缩放 90%~112%", False, 12, GRAY),
    ("  · 字体随机切换 (12%概率/词边界)", False, 12, GRAY),
    ("", False, 6, WHITE),
    ("段落级抖动", True, 14, ACCENT2),
    ("  · 左缩进随机 ±16pt", False, 12, GRAY),
    ("  · 首行缩进随机 ±10pt", False, 12, GRAY),
    ("  · 段前间距随机 ±8pt", False, 12, GRAY),
    ("  · 行间距范围 13~20pt (数学段落 16~22pt)", False, 12, GRAY),
]
add_multiline(slide, 7.4, 3.0, 4.8, 3.5, jitter_lines)

add_card(slide, 0.8, 6.3, 11.7, 0.9)
add_textbox(slide, 1.1, 6.4, 11.1, 0.7,
            "技术护城河：手写字体商用授权 + LaTeX公式渲染 + 矩阵表格排版 + 根号/上下标精确控制 = 无法快速复制的工程积累",
            font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 5)

# ────────────────────────────────────────────────────────────────
# SLIDE 6: Monetization Model
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "盈利模式", "MONETIZATION MODEL")

add_textbox(slide, 0.8, 2.1, 11.7, 0.5,
            "规则：每个新用户终身仅有 3 次免费试用 → 用完必须付费",
            font_size=16, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Products
products = [
    ("3次免费", "新用户试用", "¥0", "获客成本约¥0.15/人", ORANGE),
    ("10次", "积分包", "¥3.00", "¥0.30/题 | 毛利率 83%", ACCENT),
    ("30次", "积分包", "¥8.00", "¥0.27/题 | 毛利率 81%", ACCENT2),
    ("100次", "积分包", "¥20.00", "¥0.20/题 | 毛利率 75%", GREEN),
    ("月卡", "VIP会员", "¥15.00", "无限解题 | 毛利率 90%+", ACCENT),
    ("年卡", "VIP会员", "¥128.00", "约71折 vs 月卡", ACCENT2),
]

for i, (title, category, price, note, clr) in enumerate(products):
    x = 0.8 + i * 2.05
    add_card(slide, x, 2.8, 1.85, 3.0)
    add_textbox(slide, x + 0.15, 2.95, 1.55, 0.4, title,
                font_size=22, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 3.4, 1.55, 0.3, category,
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 3.8, 1.55, 0.5, price,
                font_size=24, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 4.5, 1.55, 1.0, note,
                font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# Conversion funnel
add_textbox(slide, 0.8, 6.1, 11.7, 0.4,
            "付费转化漏斗预计：100%注册 → 用完3次免费(85%) → 查看付费页面(40%) → 首次付费(12%) → 复购(6%)",
            font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 6)

# ────────────────────────────────────────────────────────────────
# SLIDE 7: Unit Economics
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "单位经济学", "UNIT ECONOMICS")

# Cost breakdown
add_card(slide, 0.8, 2.3, 5.8, 3.3)
add_textbox(slide, 1.1, 2.5, 5.3, 0.5, "单题成本结构 (¥0.05/题)", font_size=18, bold=True, color=ACCENT)

cost_data = [
    ("AI 视觉大模型调用", "¥0.037", "74%", ACCENT),
    ("服务器渲染 (CPU)", "¥0.005", "10%", ACCENT2),
    ("对象存储 + CDN 流量", "¥0.003", "6%", GREEN),
    ("数据库 + 其他", "¥0.005", "10%", GRAY),
]
for i, (item, cost, pct, clr) in enumerate(cost_data):
    y = 3.15 + i * 0.55
    add_textbox(slide, 1.3, y, 3.0, 0.35, item, font_size=13, color=WHITE)
    add_textbox(slide, 4.5, y, 0.8, 0.35, cost, font_size=14, bold=True, color=clr)
    add_textbox(slide, 5.3, y, 0.6, 0.35, pct, font_size=13, color=GRAY)

# Revenue
add_card(slide, 7.1, 2.3, 5.4, 3.3)
add_textbox(slide, 7.4, 2.5, 4.8, 0.5, "单用户生命周期价值 (LTV)", font_size=18, bold=True, color=ACCENT)

ltv_lines = [
    ("保守估算 (12% 付费转化)", True, 15, ACCENT2),
    ("  · 平均客单价: ¥8.00 (首次购买)", False, 13, GRAY),
    ("  · 复购率: 50% → 二次消费 ¥12.00", False, 13, GRAY),
    ("  · LTV ≈ ¥8 + 0.5 × ¥12 = ¥14.00", False, 13, GREEN),
    ("", False, 6, WHITE),
    ("获客成本 CAC", True, 15, ACCENT2),
    ("  · 3次免费试用成本: 3 × ¥0.05 = ¥0.15", False, 13, GRAY),
    ("  · 自然流量为主 + 校园推广", False, 13, GRAY),
    ("  · LTV / CAC = ¥14 / ¥0.15 ≈ 93:1", False, 13, GREEN),
]
add_multiline(slide, 7.4, 3.0, 4.8, 2.5, ltv_lines)

# Scale economics
add_card(slide, 0.8, 5.9, 11.7, 1.3)
add_textbox(slide, 1.1, 6.0, 11.1, 0.4, "规模化效应", font_size=16, bold=True, color=ACCENT)

scale_lines = [
    ("1,000 付费用户 → 月收入 ¥8,000 | 月运营成本 ¥800 | 毛利 ¥7,200/月", False, 13, WHITE),
    ("10,000 付费用户 → 月收入 ¥80,000 | 月运营成本 ¥3,000 | 毛利 ¥77,000/月", False, 13, WHITE),
    ("100,000 付费用户 → 月收入 ¥800,000 | 月运营成本 ¥15,000 | 毛利 ¥785,000/月", False, 13, GREEN),
]
add_multiline(slide, 1.1, 6.45, 11.1, 0.7, scale_lines)

add_page_number(slide, 7)

# ────────────────────────────────────────────────────────────────
# SLIDE 8: Market & Competition
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "市场与竞争格局", "MARKET & COMPETITION")

# TAM SAM SOM
add_metric_card(slide, 0.8, 2.3, "¥2800亿", "TAM 中国教育科技市场")
add_metric_card(slide, 3.9, 2.3, "¥50亿", "SAM 大学教育辅助工具")
add_metric_card(slide, 7.0, 2.3, "¥5亿", "SOM 高数解题细分市场")
add_metric_card(slide, 10.1, 2.3, "0个竞品", "手写体输出赛道")

# Competition matrix
add_card(slide, 0.8, 4.1, 11.7, 3.0)
add_textbox(slide, 1.1, 4.25, 11.1, 0.4, "竞品对比矩阵", font_size=16, bold=True, color=ACCENT)

# Table header
cols = ["", "拍照解题", "中文支持", "手写体输出", "批量处理", "价格"]
col_x = [1.1, 3.6, 5.3, 7.0, 8.7, 10.4]
for i, (hdr, x) in enumerate(zip(cols, col_x)):
    add_textbox(slide, x, 4.75, 1.5, 0.3, hdr, font_size=12, bold=True, color=ACCENT)

competitors = [
    ("高数作业杀手 (我们)", ["✅", "✅", "✅ 11款字体", "✅", "¥0.2~0.3/题"], GREEN),
    ("Photomath (Google)", ["✅", "❌ 无中文", "❌", "❌", "免费"], GRAY),
    ("Wolfram Alpha", ["❌ 仅公式", "❌", "❌", "❌", "$5/月"], GRAY),
    ("ChatGPT / Kimi", ["✅", "✅", "❌ LaTeX格式", "❌", "免费"], GRAY),
    ("作业帮 / 小猿搜题", ["✅", "✅", "❌ 打印体", "❌", "免费+会员"], GRAY),
    ("传统家教/辅导班", ["N/A", "N/A", "N/A", "N/A", "¥100~300/小时"], GRAY),
]

for i, (name, checks, clr) in enumerate(competitors):
    y = 5.2 + i * 0.38
    add_textbox(slide, 1.1, y, 2.3, 0.3, name, font_size=12, bold=True, color=clr)
    for j, check in enumerate(checks):
        add_textbox(slide, col_x[j+1], y, 1.5, 0.3, check, font_size=11, color=GRAY)

add_page_number(slide, 8)

# ────────────────────────────────────────────────────────────────
# SLIDE 9: Go-to-Market
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "增长策略", "GO-TO-MARKET STRATEGY")

strategies = [
    ("🎓", "校园裂变", "1",
     "第一批种子用户\n通过高校数学系/考研群推广\n「分享得积分」裂变机制",
     "0元获客，\n精准触达目标人群"),
    ("📱", "微信生态", "2",
     "微信搜一搜SEO优化\n公众号+视频号内容矩阵\n朋友圈分享卡片展示手写效果",
     "微信12亿月活，\n自然流量红利"),
    ("🎯", "场景渗透", "3",
     "考研备考季专项活动\n期末考试周限时促销\n与高校数学院系合作",
     "高频刚需场景，\n转化率最高时段"),
    ("🏪", "付费转化", "4",
     "3次免费→付费弹窗（试用耗尽时）\n首次付费优惠（首单5折）\n积分过期提醒→复购触发",
     "自然转化漏斗，\n不依赖广告投放"),
]

for i, (icon, title, num, desc, result) in enumerate(strategies):
    x = 0.8 + i * 3.1
    add_card(slide, x, 2.3, 2.8, 4.2)
    add_textbox(slide, x + 0.3, 2.5, 2.2, 0.5, f"{icon}  {title}",
                font_size=18, bold=True)
    # Phase number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 2.1), Inches(2.45), Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + 0.3, 3.3, 2.2, 2.0, desc, font_size=13, color=GRAY)
    add_accent_line(slide, x + 0.3, 5.35, 2.2, ACCENT2)
    add_textbox(slide, x + 0.3, 5.5, 2.2, 0.8, result, font_size=11, color=GREEN)

add_page_number(slide, 9)

# ────────────────────────────────────────────────────────────────
# SLIDE 10: Financial Projections
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "财务预测 (3年)", "FINANCIAL PROJECTIONS")

# Year columns
years = [
    ("第一年\n2026-2027", "¥120万", "1.2万", "¥1,900", [
        ("注册用户", "50,000"),
        ("付费用户", "6,000 (12%)"),
        ("年收入", "¥144,000"),
        ("年运营成本", "¥22,800"),
        ("年利润", "¥121,200"),
    ]),
    ("第二年\n2027-2028", "¥480万", "5.0万", "¥1,500", [
        ("注册用户", "200,000"),
        ("付费用户", "28,000 (14%)"),
        ("年收入", "¥672,000"),
        ("年运营成本", "¥60,000"),
        ("年利润", "¥612,000"),
    ]),
    ("第三年\n2028-2029", "¥2,000万", "20万", "¥950", [
        ("注册用户", "800,000"),
        ("付费用户", "120,000 (15%)"),
        ("年收入", "¥3,000,000"),
        ("年运营成本", "¥180,000"),
        ("年利润", "¥2,820,000"),
    ]),
]

for i, (year_label, valuation, users, cost_per, details) in enumerate(years):
    x = 0.8 + i * 4.1
    add_card(slide, x, 2.3, 3.8, 4.9)
    add_textbox(slide, x + 0.3, 2.45, 3.2, 0.7, year_label,
                font_size=15, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.15, 3.2, 0.4, f"估值: {valuation}",
                font_size=16, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, 3.5, 3.2, 0.3, f"付费用户: {users}",
                font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, x + 0.5, 3.95, 2.8, ACCENT2)
    for j, (label, value) in enumerate(details):
        add_textbox(slide, x + 0.5, 4.15 + j * 0.38, 1.8, 0.3, label,
                    font_size=12, color=GRAY)
        add_textbox(slide, x + 2.3, 4.15 + j * 0.38, 1.2, 0.3, value,
                    font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

add_page_number(slide, 10)

# ────────────────────────────────────────────────────────────────
# SLIDE 11: The Ask
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_section_title(slide, "融资需求", "THE ASK")

add_card(slide, 0.8, 2.3, 5.8, 4.5)
add_textbox(slide, 1.1, 2.5, 5.3, 0.5, "本轮融资", font_size=22, bold=True, color=ACCENT)
add_textbox(slide, 1.1, 3.2, 5.3, 0.8, "¥200万 (天使轮)", font_size=42, bold=True, color=WHITE)
add_textbox(slide, 1.1, 4.0, 5.3, 0.4, "出让 10% 股权  |  估值 ¥2000万", font_size=16, color=GRAY)

add_accent_line(slide, 1.1, 4.6, 5.3, ACCENT2)
fund_use = [
    ("资金用途", True, 16, ACCENT2),
    ("  · 核心研发团队扩充 (3人 × 12月) — ¥80万", False, 13, GRAY),
    ("  · 渲染引擎 Playwright 升级 + 字体商用授权 — ¥20万", False, 13, GRAY),
    ("  · 服务器与基础设施 (2年) — ¥30万", False, 13, GRAY),
    ("  · 市场推广与校园渠道 — ¥40万", False, 13, GRAY),
    ("  · 微信小程序审核 + 法务合规 — ¥10万", False, 13, GRAY),
    ("  · 运营储备金 — ¥20万", False, 13, GRAY),
]
add_multiline(slide, 1.1, 4.8, 5.3, 3.5, fund_use)

add_card(slide, 7.1, 2.3, 5.4, 4.5)
add_textbox(slide, 7.4, 2.5, 4.8, 0.5, "退出路径", font_size=22, bold=True, color=ACCENT)

exit_lines = [
    ("战略收购 (3-5年)", True, 16, ACCENT2),
    ("  · 教育科技巨头 (作业帮、猿辅导等)", False, 13, GRAY),
    ("  · AI平台公司 (扩充工具矩阵)", False, 13, GRAY),
    ("  · 预期估值 ¥5000万 ~ ¥2亿", False, 13, GREEN),
    ("", False, 6, WHITE),
    ("里程碑规划", True, 16, ACCENT2),
    ("  · 6个月: MVP上线，5000注册用户", False, 13, GRAY),
    ("  · 12个月: 5万用户，月收入 ¥1.2万", False, 13, GRAY),
    ("  · 18个月: 20万用户，月收入 ¥5万+", False, 13, GRAY),
    ("  · 24个月: 启动 A 轮融资 (¥800万)", False, 13, GREEN),
]
add_multiline(slide, 7.4, 3.0, 4.8, 3.5, exit_lines)

add_textbox(slide, 0.8, 6.6, 11.7, 0.5,
            "联系方式：[请填写]  |  邮箱：[请填写]  |  微信：[请填写]",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 11)

# ────────────────────────────────────────────────────────────────
# SLIDE 12: Thank You
# ────────────────────────────────────────────────────────────────
slide = new_slide()
add_accent_line(slide, 4.5, 3.0, 4.3, ACCENT)
add_textbox(slide, 2.5, 2.0, 8.3, 0.8, "谢谢",
            font_size=48, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.5, 3.3, 8.3, 0.6, "高数作业杀手",
            font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2.5, 4.2, 8.3, 1.0,
            "让 AI 替你写作业，让手写体骗过老师\n\n期待与您合作",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 12)

# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "高数作业杀手_融资路演.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
